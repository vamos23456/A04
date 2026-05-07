import os
from io import BytesIO
from typing import List, Optional

from docx import Document
from fastapi import APIRouter, Depends, File, HTTPException, UploadFile
from pydantic import BaseModel
from pptx import Presentation
from sqlalchemy.orm import Session

from auth import get_current_user, get_optional_current_user
from database import KnowledgeDocument, User, get_db
from rag_service import rag_service

router = APIRouter(prefix="/api/knowledge", tags=["知识库"])

DEFAULT_KNOWLEDGE_FILES = [
    "../../knowledge/初中物理知识点总结大全.txt",
    "../../knowledge/初中数学知识点总结大全.txt",
]


class SearchRequest(BaseModel):
    query: str
    top_k: int = 3


class KnowledgeDocumentCreate(BaseModel):
    title: str
    content: str
    source_name: Optional[str] = None


class KnowledgeDocumentResponse(BaseModel):
    id: int
    title: str
    source_name: Optional[str] = None
    chunk_count: int
    is_system: bool
    owner_scope: str
    created_at: str
    updated_at: str


def _extract_text_from_file(file_name: str, content: bytes) -> str:
    extension = os.path.splitext(file_name)[1].lower()

    if extension in {".txt", ".md", ".csv"}:
        return content.decode("utf-8", errors="ignore").strip()

    if extension == ".docx":
        document = Document(BytesIO(content))
        return "\n".join(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()).strip()

    if extension == ".pptx":
        presentation = Presentation(BytesIO(content))
        texts: List[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts).strip()

    raise HTTPException(status_code=400, detail="暂不支持该文件类型，请上传 txt、md、csv、docx 或 pptx 文件")


def _seed_default_knowledge(db: Session):
    """把默认知识库文档写入数据库，并建立公共向量索引。"""
    total_files = 0
    total_chunks = 0
    base_dir = os.path.dirname(os.path.abspath(__file__))

    for relative_path in DEFAULT_KNOWLEDGE_FILES:
        absolute_path = os.path.join(base_dir, relative_path)
        source_name = os.path.basename(absolute_path)

        if not os.path.exists(absolute_path):
            continue

        exists = db.query(KnowledgeDocument).filter(
            KnowledgeDocument.user_id.is_(None),
            KnowledgeDocument.is_system.is_(True),
            KnowledgeDocument.source_name == source_name
        ).first()
        if exists:
            continue

        with open(absolute_path, "r", encoding="utf-8") as file:
            content = file.read().strip()

        if not content:
            continue

        record = KnowledgeDocument(
            title=os.path.splitext(source_name)[0],
            source_name=source_name,
            content=content,
            is_system=True,
        )
        db.add(record)
        db.commit()
        db.refresh(record)

        index_result = rag_service.index_document(
            document_id=record.id,
            text=record.content,
            title=record.title,
            source_name=record.source_name,
            is_system=True,
        )
        record.chunk_count = index_result["total_chunks"]
        db.commit()

        total_files += 1
        total_chunks += record.chunk_count

    return {"total_files": total_files, "total_chunks": total_chunks}


@router.post("/init")
def init_knowledge(db: Session = Depends(get_db)):
    """初始化默认公共知识库。"""
    try:
        result = _seed_default_knowledge(db)
        return {
            "message": "默认知识库初始化完成",
            "total_files": result["total_files"],
            "total_chunks": result["total_chunks"],
        }
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"初始化失败：{str(exc)}")


@router.post("/search")
def search_knowledge(
    req: SearchRequest,
    current_user: Optional[User] = Depends(get_optional_current_user),
    db: Session = Depends(get_db),
):
    """检索公共知识库，并在登录时叠加当前用户私有知识库。"""
    try:
        _seed_default_knowledge(db)
        results = rag_service.search(req.query, req.top_k, user_id=current_user.id if current_user else None)
        return {
            "query": req.query,
            "documents": results["documents"],
            "metadatas": results["metadatas"],
        }
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"检索失败：{str(exc)}")


@router.get("/stats")
def get_knowledge_stats(
    current_user: Optional[User] = Depends(get_optional_current_user),
    db: Session = Depends(get_db),
):
    """获取公共知识库和当前用户私有知识库统计。"""
    try:
        _seed_default_knowledge(db)
        public_docs = db.query(KnowledgeDocument).filter(KnowledgeDocument.user_id.is_(None)).all()
        user_docs = []
        if current_user:
            user_docs = db.query(KnowledgeDocument).filter(KnowledgeDocument.user_id == current_user.id).all()

        return rag_service.get_stats(
            public_documents=len(public_docs),
            public_chunks=sum(doc.chunk_count for doc in public_docs),
            user_documents=len(user_docs),
            user_chunks=sum(doc.chunk_count for doc in user_docs),
        )
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"获取统计失败：{str(exc)}")


@router.get("/documents", response_model=List[KnowledgeDocumentResponse])
def list_documents(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """列出默认知识库和当前用户的私有知识库。"""
    _seed_default_knowledge(db)
    documents = db.query(KnowledgeDocument).filter(
        (KnowledgeDocument.user_id.is_(None)) | (KnowledgeDocument.user_id == current_user.id)
    ).order_by(KnowledgeDocument.is_system.desc(), KnowledgeDocument.created_at.desc()).all()

    return [
        {
            "id": doc.id,
            "title": doc.title,
            "source_name": doc.source_name,
            "chunk_count": doc.chunk_count,
            "is_system": doc.is_system,
            "owner_scope": "public" if doc.user_id is None else "private",
            "created_at": doc.created_at.isoformat(),
            "updated_at": doc.updated_at.isoformat(),
        }
        for doc in documents
    ]


@router.post("/documents", response_model=KnowledgeDocumentResponse)
def create_document(
    payload: KnowledgeDocumentCreate,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """向当前用户的私有知识库添加一条文档。"""
    content = payload.content.strip()
    if not content:
        raise HTTPException(status_code=400, detail="知识内容不能为空")

    document = KnowledgeDocument(
        user_id=current_user.id,
        title=payload.title.strip(),
        source_name=(payload.source_name or payload.title).strip(),
        content=content,
        is_system=False,
    )
    db.add(document)
    db.commit()
    db.refresh(document)

    try:
        result = rag_service.index_document(
            document_id=document.id,
            text=document.content,
            title=document.title,
            source_name=document.source_name,
            user_id=current_user.id,
        )
        document.chunk_count = result["total_chunks"]
        db.commit()
        db.refresh(document)
    except Exception as exc:
        db.delete(document)
        db.commit()
        raise HTTPException(status_code=500, detail=f"索引知识文档失败：{str(exc)}")

    return {
        "id": document.id,
        "title": document.title,
        "source_name": document.source_name,
        "chunk_count": document.chunk_count,
        "is_system": document.is_system,
        "owner_scope": "private",
        "created_at": document.created_at.isoformat(),
        "updated_at": document.updated_at.isoformat(),
    }


@router.post("/documents/upload", response_model=KnowledgeDocumentResponse)
async def upload_document(
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """上传文档并写入当前用户私有知识库。"""
    file_bytes = await file.read()
    if not file_bytes:
        raise HTTPException(status_code=400, detail="上传文件为空")

    text = _extract_text_from_file(file.filename or "未命名文件", file_bytes)
    if not text:
        raise HTTPException(status_code=400, detail="文件中没有可提取的文字内容")

    title = os.path.splitext(file.filename or "未命名文件")[0]
    document = KnowledgeDocument(
        user_id=current_user.id,
        title=title,
        source_name=file.filename,
        content=text,
        is_system=False,
    )
    db.add(document)
    db.commit()
    db.refresh(document)

    try:
        result = rag_service.index_document(
            document_id=document.id,
            text=document.content,
            title=document.title,
            source_name=document.source_name,
            user_id=current_user.id,
        )
        document.chunk_count = result["total_chunks"]
        db.commit()
        db.refresh(document)
    except Exception as exc:
        db.delete(document)
        db.commit()
        raise HTTPException(status_code=500, detail=f"上传入库失败：{str(exc)}")

    return {
        "id": document.id,
        "title": document.title,
        "source_name": document.source_name,
        "chunk_count": document.chunk_count,
        "is_system": document.is_system,
        "owner_scope": "private",
        "created_at": document.created_at.isoformat(),
        "updated_at": document.updated_at.isoformat(),
    }


@router.delete("/documents/{document_id}")
def delete_document(
    document_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """删除当前用户自己的知识文档。"""
    document = db.query(KnowledgeDocument).filter(
        KnowledgeDocument.id == document_id,
        KnowledgeDocument.user_id == current_user.id
    ).first()

    if not document:
        raise HTTPException(status_code=404, detail="知识文档不存在")

    rag_service.delete_document(document_id=document.id, user_id=current_user.id)
    db.delete(document)
    db.commit()

    return {"message": "知识文档已删除"}
