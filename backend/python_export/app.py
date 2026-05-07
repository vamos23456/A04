from io import BytesIO
from typing import List, Optional
from urllib.parse import quote
from urllib.request import urlopen
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from fastapi import FastAPI
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response
from pydantic import BaseModel, Field
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, RGBColor as DocxRGB
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt as PPt
from pptx.enum.text import PP_ALIGN

# 导入新增的路由
from auth_routes import router as auth_router
from knowledge_routes import router as knowledge_router
from project_routes import router as project_router
from database import init_db

app = FastAPI(title="AI Teaching Assistant API")

# 初始化数据库
init_db()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 注册路由
app.include_router(auth_router)
app.include_router(knowledge_router)
app.include_router(project_router)


# ── Models ────────────────────────────────────────────────────────────────────

class LessonProcessStep(BaseModel):
    stage: str = ""
    duration: str = ""
    activities: List[str] = Field(default_factory=list)


class LessonPlan(BaseModel):
    title: str = "未命名教案"
    objectives: List[str] = Field(default_factory=list)
    process: List[LessonProcessStep] = Field(default_factory=list)
    methods: List[str] = Field(default_factory=list)
    homework: List[str] = Field(default_factory=list)


class PptSlide(BaseModel):
    layout: str = "content"
    title: Optional[str] = ""
    subtitle: Optional[str] = ""
    body: List[str] = Field(default_factory=list)
    imageKeyword: Optional[str] = ""
    accent: Optional[str] = "6366F1"
    note: Optional[str] = ""


class ExportDocxRequest(BaseModel):
    word: LessonPlan


class ExportPptRequest(BaseModel):
    title: str = "PPT课件"
    slides: List[PptSlide] = Field(default_factory=list)


# ── Helpers ───────────────────────────────────────────────────────────────────

FONT_CN = "微软雅黑"

def _set_font_cn(run, font_name: str = FONT_CN):
    """同时设置西文和东亚（中文）字体，避免 Word 回退到宋体。"""
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.insert(0, rFonts)
    rFonts.set(qn("w:eastAsia"), font_name)
    rFonts.set(qn("w:hint"), "eastAsia")


def _set_style_font_cn(style, font_name: str = FONT_CN):
    """给 style 的 rPr 设置东亚字体，影响该样式下所有段落。"""
    style.font.name = font_name
    rPr = style.element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.insert(0, rFonts)
    rFonts.set(qn("w:eastAsia"), font_name)

def _hex(color: Optional[str], default: str = "6366F1") -> str:
    if not color:
        return default
    c = color.replace("#", "").strip().upper()
    return c if len(c) == 6 else default


def _rgb(color: Optional[str], default: str = "6366F1") -> RGBColor:
    h = _hex(color, default)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _safe_name(name: str, fallback: str) -> str:
    import re
    base = (name or "").strip() or fallback
    return re.sub(r'[\\/:*?"<>|]', "_", base)


def _blank_slide(prs: Presentation):
    """Add a fully blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]  # index 6 = Blank
    return prs.slides.add_slide(blank_layout)


def _add_rect(slide, x, y, w, h, hex_color: str):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(
        int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    )
    shape.line.fill.background()
    return shape


def _add_textbox(slide, x, y, w, h, text: str, *,
                 size=18, bold=False, color="1A1A1A",
                 align=PP_ALIGN.LEFT, valign="top",
                 font_name="Microsoft YaHei", wrap=True):
    from pptx.util import Inches
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = _rgb(color)
    return txBox


def _add_bullets(slide, x, y, w, h, items: List[str], *,
                 size=18, color="333333", font_name="Microsoft YaHei"):
    from pptx.util import Inches
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = PPt(6)
        for run in p.runs:
            run.font.size = PPt(size)
            run.font.name = font_name
            run.font.color.rgb = _rgb(color)
    return txBox


def _fetch_image(keyword: str) -> Optional[bytes]:
    """从 picsum.photos 获取图片，超时时间缩短为 3 秒"""
    try:
        url = f"https://picsum.photos/seed/{keyword or 'education'}/800/600"
        with urlopen(url, timeout=3) as r:
            return r.read()
    except Exception:
        return None


# ── Slide renderers ───────────────────────────────────────────────────────────

W, H = 13.333, 7.5  # slide dimensions in inches


def _render_cover(prs, slide_data: PptSlide):
    s = _blank_slide(prs)
    acc = _hex(slide_data.accent)
    # full background
    _add_rect(s, 0, 0, W, H, acc)
    # title
    _add_textbox(s, 1, 2.2, W - 2, 1.5,
                 slide_data.title or "",
                 size=44, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
    # subtitle
    if slide_data.subtitle:
        _add_textbox(s, 1, 3.9, W - 2, 0.9,
                     slide_data.subtitle,
                     size=24, color="FFFFFF", align=PP_ALIGN.CENTER)
    _add_note(s, slide_data.note)


def _render_content(prs, slide_data: PptSlide):
    s = _blank_slide(prs)
    acc = _hex(slide_data.accent)
    # top bar
    _add_rect(s, 0, 0, W, 0.15, acc)
    # title
    _add_textbox(s, 0.6, 0.25, W - 1.2, 1.0,
                 slide_data.title or "",
                 size=32, bold=True, color="1A1A1A")
    # vertical accent line
    _add_rect(s, 0.6, 1.4, 0.07, H - 1.8, acc)
    # bullets
    if slide_data.body:
        _add_bullets(s, 0.9, 1.45, W - 1.5, H - 1.9,
                     slide_data.body, size=20)
    _add_note(s, slide_data.note)


def _render_two_col(prs, slide_data: PptSlide):
    s = _blank_slide(prs)
    acc = _hex(slide_data.accent)
    _add_rect(s, 0, 0, W, 0.15, acc)
    _add_textbox(s, 0.6, 0.25, W - 1.2, 0.9,
                 slide_data.title or "",
                 size=30, bold=True, color="1A1A1A")
    body = slide_data.body or []
    mid = (len(body) + 1) // 2
    left, right = body[:mid], body[mid:]
    col_y, col_h = 1.3, H - 1.7
    col_w = (W - 1.2) / 2 - 0.15
    # left card bg
    _add_rect(s, 0.5, col_y, col_w, col_h, "F5F5F5")
    if left:
        _add_bullets(s, 0.7, col_y + 0.2, col_w - 0.4, col_h - 0.4,
                     left, size=18)
    # right card bg
    rx = 0.5 + col_w + 0.3
    _add_rect(s, rx, col_y, col_w, col_h, "F5F5F5")
    if right:
        _add_bullets(s, rx + 0.2, col_y + 0.2, col_w - 0.4, col_h - 0.4,
                     right, size=18)
    _add_note(s, slide_data.note)


def _render_image_text(prs, slide_data: PptSlide):
    s = _blank_slide(prs)
    acc = _hex(slide_data.accent)
    _add_rect(s, 0, 0, W, 0.15, acc)
    _add_textbox(s, 0.6, 0.25, W - 1.2, 0.9,
                 slide_data.title or "",
                 size=30, bold=True, color="1A1A1A")
    img_bytes = _fetch_image(slide_data.imageKeyword or "education")
    if img_bytes:
        img_buf = BytesIO(img_bytes)
        s.shapes.add_picture(img_buf, Inches(0.5), Inches(1.3),
                             Inches(5.5), Inches(H - 1.7))
    else:
        # fallback color block
        _add_rect(s, 0.5, 1.3, 5.5, H - 1.7, acc)
        _add_textbox(s, 0.5, H / 2 - 0.3, 5.5, 0.6,
                     slide_data.imageKeyword or "",
                     size=16, color="FFFFFF", align=PP_ALIGN.CENTER)
    if slide_data.body:
        _add_bullets(s, 6.3, 1.3, W - 6.8, H - 1.7,
                     slide_data.body, size=19)
    _add_note(s, slide_data.note)


def _render_quote(prs, slide_data: PptSlide):
    s = _blank_slide(prs)
    acc = _hex(slide_data.accent)
    _add_rect(s, 0, 0, W, H, acc)
    # big quote mark
    _add_textbox(s, 0.5, 0.2, 3, 2, "\u201C",
                 size=140, color="FFFFFF", align=PP_ALIGN.LEFT)
    _add_textbox(s, 1.2, 2.0, W - 2.4, 2.5,
                 slide_data.title or "",
                 size=30, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
    if slide_data.subtitle:
        _add_textbox(s, 1.2, 4.7, W - 2.4, 0.7,
                     f"\u2014\u2014 {slide_data.subtitle}",
                     size=18, color="FFFFFF", align=PP_ALIGN.RIGHT)
    _add_note(s, slide_data.note)


def _render_end(prs, slide_data: PptSlide):
    s = _blank_slide(prs)
    acc = _hex(slide_data.accent)
    _add_rect(s, 0, 0, W, H, "1A1A1A")
    # accent line
    _add_rect(s, W / 2 - 1.5, H / 2 - 0.3, 3, 0.08, acc)
    _add_textbox(s, 1, H / 2 - 0.1, W - 2, 1.3,
                 slide_data.title or "谢谢",
                 size=48, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
    if slide_data.subtitle:
        _add_textbox(s, 1, H / 2 + 1.3, W - 2, 0.7,
                     slide_data.subtitle,
                     size=20, color="AAAAAA", align=PP_ALIGN.CENTER)
    _add_note(s, slide_data.note)


def _add_note(slide, note: Optional[str]):
    if note:
        slide.notes_slide.notes_text_frame.text = note


RENDERERS = {
    "cover": _render_cover,
    "content": _render_content,
    "two-col": _render_two_col,
    "image-text": _render_image_text,
    "quote": _render_quote,
    "end": _render_end,
}


# ── Routes ────────────────────────────────────────────────────────────────────

@app.get("/health")
def health():
    return {"ok": True}


@app.post("/export/docx")
def export_docx(payload: ExportDocxRequest) -> Response:
    word = payload.word
    doc = Document()

    # 全局字体：同时设置西文和东亚字体
    for sname in ("Normal", "Heading 1", "Heading 2", "List Bullet", "List Bullet 2"):
        try:
            _set_style_font_cn(doc.styles[sname])
            doc.styles[sname].font.size = Pt(12)
        except KeyError:
            pass

    title_p = doc.add_heading(word.title or "未命名教案", level=0)
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in title_p.runs:
        _set_font_cn(run)

    doc.add_heading("教学目标", level=1)
    for item in word.objectives:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(item)
        _set_font_cn(run)

    doc.add_heading("教学过程", level=1)
    for step in word.process:
        p = doc.add_paragraph()
        run = p.add_run(step.stage or "未命名环节")
        run.bold = True
        run.font.size = Pt(13)
        _set_font_cn(run)
        if step.duration:
            r2 = p.add_run(f"  （{step.duration}）")
            r2.font.size = Pt(11)
            r2.font.color.rgb = DocxRGB(0x99, 0x99, 0x99)
            _set_font_cn(r2)
        for activity in step.activities:
            p2 = doc.add_paragraph(style="List Bullet 2")
            run2 = p2.add_run(activity)
            _set_font_cn(run2)

    doc.add_heading("教学方法", level=1)
    if word.methods:
        p = doc.add_paragraph()
        run = p.add_run("、".join(word.methods))
        _set_font_cn(run)

    doc.add_heading("课后作业", level=1)
    for hw in word.homework:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(hw)
        _set_font_cn(run)

    buf = BytesIO()
    doc.save(buf)
    name = _safe_name(word.title or "教案", "教案") + ".docx"
    return Response(
        content=buf.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(name)}"},
    )


@app.post("/export/pptx")
def export_pptx(payload: ExportPptRequest) -> Response:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    for slide_data in payload.slides:
        renderer = RENDERERS.get(slide_data.layout, _render_content)
        renderer(prs, slide_data)

    if not payload.slides:
        _render_end(prs, PptSlide(title=payload.title or "PPT课件"))

    buf = BytesIO()
    prs.save(buf)
    name = _safe_name(payload.title or "PPT课件", "PPT课件") + ".pptx"
    return Response(
        content=buf.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(name)}"},
    )
